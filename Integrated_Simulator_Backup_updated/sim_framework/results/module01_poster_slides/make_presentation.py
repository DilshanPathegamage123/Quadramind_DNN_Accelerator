#!/usr/bin/env python3
"""Build the Module 01 presentation update deck.

These are ADDITIONS and CORRECTIONS to the existing final-evaluation deck
(slides 13-27).  The deck is deliberately short: three new slides that close
gaps an evaluator is likely to probe, one corrected verification slide, and a
closing summary.

Styling follows the existing deck: blue header band, white bold title,
dark-blue body text on white.

Run:  ~/simenv/bin/python results/module01_poster_slides/make_presentation.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Module01_Presentation_Update.pptx"

BLUE = RGBColor(0x0B, 0x4D, 0xA2)      # header band / accents
NAVY = RGBColor(0x12, 0x35, 0x5B)      # body text
TEAL = RGBColor(0x0F, 0x7B, 0x6C)      # good / measured
ORANGE = RGBColor(0xC1, 0x7D, 0x11)    # caution / modelled
GREY = RGBColor(0x5A, 0x66, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE8, 0xEF, 0xF8)

W, H = Inches(13.333), Inches(7.5)


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, title: str, band_h: float = 1.15) -> None:
    """Blue band across the top with the slide title in white."""
    box = slide.shapes.add_shape(1, 0, 0, W, Inches(band_h))  # 1 = rectangle
    box.fill.solid()
    box.fill.fore_color.rgb = BLUE
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left, tf.margin_right = Inches(0.5), Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    p.runs[0].font.size = Pt(30)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def bullets(tf, items, size=18, color=NAVY, space=10) -> None:
    """items: list of (text, indent_level) or plain strings."""
    for i, item in enumerate(items):
        text, lvl = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("- " if lvl else "") + text
        p.level = lvl
        p.space_after = Pt(space)
        for r in p.runs:
            r.font.size = Pt(size - 2 * lvl)
            r.font.color.rgb = color


def chip(slide, left, top, width, height, title, body, accent=BLUE,
         title_size=16, body_size=13):
    """A rounded light card with a bold heading and a short body.

    The text lives in a separate textbox laid over the shape rather than in the
    shape's own text frame: autoshape text inherits theme styling and gets
    clipped to the shape bounds, which silently mangles longer captions.
    A plain textbox neither clips nor inherits.
    """
    card = slide.shapes.add_shape(5, left, top, width, height)  # 5 = rounded rect
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = accent
    card.line.width = Pt(1.25)

    tf = textbox(slide, left, top + Inches(0.12), width, height - Inches(0.2))
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(title_size)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = accent
    if body:  # callers pass "" when they fill the card themselves
        p2 = tf.add_paragraph()
        p2.text = body
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        # "\n" splits the paragraph into several runs; style every one of them
        # or the trailing lines fall back to the 18 pt default.
        for r in p2.runs:
            r.font.size = Pt(body_size)
            r.font.color.rgb = NAVY
    return tf


def table(slide, rows_data, left, top, width, height,
          col_widths=None, size=13, header_size=13):
    rows, cols = len(rows_data), len(rows_data[0])
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            # "\n" in a cell creates extra PARAGRAPHS (not runs), so style
            # every paragraph or the wrapped lines revert to the 18 pt default.
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(header_size if r == 0 else size)
                    run.font.bold = (r == 0)
                    run.font.color.rgb = WHITE if r == 0 else NAVY
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r == 0 else (
                WHITE if r % 2 else LIGHT)
    return tbl


def caption(slide, text, top, size=13, color=GREY):
    tf = textbox(slide, Inches(0.6), top, Inches(12.1), Inches(0.6))
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(size)
    p.runs[0].font.italic = True
    p.runs[0].font.color.rgb = color


# ===========================================================================
# Slide 1 - how the whole work fits together (the flow)
# ===========================================================================
def slide_flow(prs):
    s = blank(prs)
    header(s, "How the Work Fits Together")

    # Keep each body to two short lines - the cards are narrow and any more
    # text overflows the rounded box when rendered.
    steps = [
        ("1. BUILD", "27 selectable\nRTL configurations", BLUE),
        ("2. VERIFY", "RTL vs TensorFlow\ngolden reference", TEAL),
        ("3. MEASURE", "Beats, bursts and\ncycles per knob", TEAL),
        ("4. MODEL", "Derive the\ncycle equation", BLUE),
        ("5. CHECK", "Validate back\nagainst RTL", TEAL),
        ("6. APPLY", "Score 27 combos\nin milliseconds", ORANGE),
    ]
    x, y, cw, gap = Inches(0.42), Inches(1.6), Inches(1.95), Inches(0.14)
    for i, (t, b, col) in enumerate(steps):
        chip(s, x + i * (cw + gap), y, cw, Inches(1.65), t, b, col,
             title_size=15, body_size=12)
        if i < len(steps) - 1:
            ar = s.shapes.add_shape(13, x + i * (cw + gap) + cw + Inches(0.01),
                                    y + Inches(0.68), gap - Inches(0.02),
                                    Inches(0.3))  # 13 = right arrow
            ar.fill.solid()
            ar.fill.fore_color.rgb = GREY
            ar.line.fill.background()

    tf = textbox(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(1.9))
    bullets(tf, [
        "Steps 2 and 3 are different experiments. Step 2 asks \"is the answer "
        "correct?\"  Step 3 asks \"how much traffic and time did it cost?\"",
        "Steps 4 and 5 form a loop: the model is built from measured hardware, "
        "then tested back against it.",
        "Only step 6 runs without RTL - that is the point of the analytical "
        "framework.",
    ], size=17)

    box = s.shapes.add_shape(5, Inches(0.6), Inches(6.05), Inches(12.1),
                             Inches(0.85))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    p = box.text_frame.paragraphs[0]
    p.text = ("My contribution is step 1 (the measurable knobs) "
              "and steps 4-6 (the analytical framework).")
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(18)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE


# ===========================================================================
# Slide 2 - measured vs modelled vs assumed
# ===========================================================================
def slide_provenance(prs):
    s = blank(prs)
    header(s, "Where Every Number Comes From")

    caption(s, "Three kinds of numbers appear in this work. They are labelled "
               "everywhere, and they are not mixed.", Inches(1.3))

    rows = [
        ["Type", "What it is", "Where it appears", "How strong"],
        ["MEASURED\n(RTL)",
         "Real Verilator runs: cycles, AXI beats,\nread bursts, output values",
         "Golden check - 22 verified\nconfigurations",
         "Hard evidence"],
        ["MODEL",
         "Closed-form predictions from my\nanalytical framework",
         "The chooser; the edge / cloud\ndesign-space study",
         "Validated against\nthe measured set"],
        ["ASSUMED",
         "Literature constants\n(DRAM 560 pJ/byte, MAC 0.5 pJ)",
         "Energy figures",
         "Cited, listed in\nthe appendix"],
    ]
    table(s, rows, Inches(0.6), Inches(1.95), Inches(12.1), Inches(2.9),
          col_widths=[Inches(1.7), Inches(4.0), Inches(3.6), Inches(2.8)],
          size=13)

    tf = textbox(s, Inches(0.6), Inches(5.15), Inches(12.1), Inches(1.9))
    bullets(tf, [
        "The large-network results (ResNet-50, BERT-Large, DLRM, VGG-16) are "
        "MODEL, never RTL. Every such chart is captioned that way.",
        "Energy is a model quantity: RTL gives no power number. But its main "
        "input - off-chip bytes - is validated exactly, so relative energy "
        "comparisons inherit that accuracy.",
    ], size=17)


# ===========================================================================
# Slide 3 - why large networks are not simulated in RTL
# ===========================================================================
def slide_why_no_big_rtl(prs):
    s = blank(prs)
    header(s, "Why Large Networks Are Not Run in RTL")

    tf = textbox(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.9))
    p = tf.paragraphs[0]
    p.text = ("RTL simulation runs cycle by cycle - about 10,000 simulated "
              "cycles per second. A faster machine changes this by 2-5x, "
              "not by orders of magnitude.")
    p.runs[0].font.size = Pt(17)
    p.runs[0].font.color.rgb = NAVY

    rows = [
        ["Task", "Cycles", "RTL time"],
        ["tiny_cnn layer_00, one configuration", "5,808", "1.3 s"],
        ["One BERT-Large layer, one configuration", "~95 M", "~2.6 hours"],
        ["The same layer x 27 combinations", "", "~3 months"],
        ["56 evaluation layers x 27 combinations", "", "years of compute"],
    ]
    table(s, rows, Inches(1.5), Inches(2.4), Inches(10.3), Inches(2.2),
          col_widths=[Inches(6.0), Inches(2.0), Inches(2.3)], size=15)

    box = s.shapes.add_shape(5, Inches(0.6), Inches(4.75), Inches(12.1),
                             Inches(2.45))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = BLUE
    box.line.width = Pt(1.5)
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = Inches(0.3)
    tf2.margin_top = Inches(0.2)
    p1 = tf2.paragraphs[0]
    p1.text = ("This is not a limitation of my setup - RTL simulation does not "
               "scale to design-space exploration, by construction.")
    p1.runs[0].font.size = Pt(18)
    p1.runs[0].font.bold = True
    p1.runs[0].font.color.rgb = NAVY
    p2 = tf2.add_paragraph()
    p2.text = ("That is exactly why the analytical framework is needed. It "
               "scores 27 combinations in about 1 millisecond, against a "
               ">= 147 s lower bound for the equivalent RTL sweep on the "
               "smallest layer - roughly 100,000x faster.")
    p2.runs[0].font.size = Pt(17)
    p2.runs[0].font.color.rgb = NAVY
    p3 = tf2.add_paragraph()
    p3.text = ("And it is not needed: weights never affect cycles or traffic, "
               "so only workload dimensions matter. Bigger networks repeat the "
               "same validated computation more times.")
    p3.runs[0].font.size = Pt(16)
    p3.runs[0].font.color.rgb = TEAL
    p3.space_before = Pt(6)


# ===========================================================================
# Slide 4 - corrected verification slide (replaces the existing one)
# ===========================================================================
def slide_verification_fixed(prs):
    s = blank(prs)
    header(s, "Verification - RTL Output vs TensorFlow Golden")

    stats = [("22 / 22", "configurations\nverified"),
             ("0.0499 %", "worst output\nerror"),
             ("5 %", "tolerance\ngate"),
             ("~100x", "safety margin\ninside tolerance")]
    x, cw, gap = Inches(0.7), Inches(2.85), Inches(0.25)
    for i, (big, lab) in enumerate(stats):
        card = s.shapes.add_shape(5, x + i * (cw + gap), Inches(1.45), cw,
                                  Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.fill.background()
        # Text overlaid, not inside the shape - see chip() for why.
        tf = textbox(s, x + i * (cw + gap), Inches(1.6), cw, Inches(1.5))
        p = tf.paragraphs[0]
        p.text = big
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(28)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = BLUE
        p2 = tf.add_paragraph()
        p2.text = lab
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(2)
        for r in p2.runs:          # "\n" -> multiple runs; style them all
            r.font.size = Pt(12)
            r.font.color.rgb = NAVY

    caption(s, "Sample evidence - 2 of the 22 verified configurations",
            Inches(3.32), size=14, color=NAVY)

    rows = [["Configuration", "Measured (RTL)", "Golden (TF)", "Error (% FS)",
             "Result"],
            ["tiny_cnn L0 - OS-ChM-Multicast", "0.709772", "0.710062", "0.01 %",
             "PASS"],
            ["mnist_cnn L2 - OS-ChM-Multicast", "-2.005136", "-2.006291",
             "0.05 %", "PASS"]]
    table(s, rows, Inches(0.7), Inches(3.7), Inches(11.9), Inches(1.2),
          col_widths=[Inches(4.3), Inches(2.0), Inches(1.9), Inches(1.8),
                      Inches(1.9)], size=14)

    box = s.shapes.add_shape(5, Inches(0.7), Inches(5.2), Inches(11.9),
                             Inches(1.55))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    tf3 = box.text_frame
    tf3.word_wrap = True
    tf3.margin_left = tf3.margin_right = Inches(0.3)
    tf3.margin_top = Inches(0.15)
    p1 = tf3.paragraphs[0]
    p1.text = ("Every configuration computes the mathematically correct answer, "
               "with roughly 100x margin inside tolerance.")
    p1.alignment = PP_ALIGN.CENTER
    p1.runs[0].font.size = Pt(18)
    p1.runs[0].font.bold = True
    p1.runs[0].font.color.rgb = WHITE
    p2 = tf3.add_paragraph()
    p2.text = ("Verified across all 3 stationary schemes, all 3 memory layouts, "
               "all 3 casting schemes, both memory backends (STAMP / PAGED) and "
               "three array geometries (8x8, 8x2, 8x1).")
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(14)
    p2.runs[0].font.color.rgb = WHITE

    caption(s, "CORRECTION vs the earlier deck: this said \"29 workloads\". 29 is "
               "the raw file count including 7 debug probe runs (worst error "
               "106 %). The verified set is 22 CONFIGURATIONS built from 5 "
               "layer shapes - and \"configuration\" is the correct word, not "
               "\"workload\".", Inches(6.85), size=12, color=ORANGE)


# ===========================================================================
# Slide 5 - what generalises, and the honest limit
# ===========================================================================
def slide_generalisation(prs):
    s = blank(prs)
    header(s, "What Generalises - and What Does Not")

    tf = chip(s, Inches(0.6), Inches(1.5), Inches(5.95), Inches(4.0),
              "WHAT GENERALISES", "", TEAL)
    for txt in [
        "The framework decides from workload GEOMETRY, not from the trained "
        "model. Weight values never affect cycles or traffic.",
        "All three stationary schemes, all three layouts, all three casting "
        "schemes and three array sizes are covered by measurement.",
        "The traffic model has NO fitted constants and matches measured beats "
        "exactly on 17 of 19 runs - including array widths it was never fitted to.",
        "K and output size enter as loop counts, so bigger networks repeat a "
        "validated computation more times.",
    ]:
        p = tf.add_paragraph()
        p.text = "- " + txt
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(9)
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.color.rgb = NAVY

    tf2 = chip(s, Inches(6.78), Inches(1.5), Inches(5.95), Inches(4.0),
               "HONEST LIMITS", "", ORANGE)
    for txt in [
        "The RTL sweep varies ONE knob at a time. Each knob's own effect is "
        "measured; knob interactions are composed analytically, not measured.",
        "The compute term is interpolated for layer shapes with no measured run "
        "- it cannot change the layout or casting choice, but it can affect the "
        "stationary-scheme choice.",
        "Energy is a model quantity; RTL produces no power number.",
        "Depthwise / grouped convolution is outside the current model.",
    ]:
        p = tf2.add_paragraph()
        p.text = "- " + txt
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(9)
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.color.rgb = NAVY

    box = s.shapes.add_shape(5, Inches(0.6), Inches(5.75), Inches(12.13),
                             Inches(1.05))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    tfb = box.text_frame
    tfb.word_wrap = True
    tfb.margin_left = tfb.margin_right = Inches(0.3)
    p = tfb.paragraphs[0]
    p.text = ("Claim made: validated for workload-based behaviour across every "
              "axis the framework selects over - not for every DNN model.")
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(17)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE


# ===========================================================================
# Slide 6 - closing summary
# ===========================================================================
def slide_summary(prs):
    s = blank(prs)
    header(s, "Module 01 - Summary")

    tf = textbox(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.6))
    p = tf.paragraphs[0]
    p.text = "Contribution"
    p.runs[0].font.size = Pt(21)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = BLUE

    tf2 = textbox(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(1.5))
    bullets(tf2, [
        "Made stationary scheme, memory layout and casting scheme "
        "simultaneously selectable and measurable in cycle-accurate RTL - "
        "27 configurations from one module. No existing tool allows this.",
        "Distilled the measured behaviour into an analytical framework that "
        "picks the best of the 27 in milliseconds.",
    ], size=17)

    tf3 = textbox(s, Inches(0.6), Inches(3.5), Inches(12.1), Inches(0.5))
    p = tf3.paragraphs[0]
    p.text = "Key findings"
    p.runs[0].font.size = Pt(21)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = BLUE

    rows = [["Finding", "Evidence"],
            ["Casting is the biggest lever on off-chip traffic",
             "1,836 -> 11,664 -> 20,736 beats (11.3x)"],
            ["Layout changes REQUESTS, not data volume",
             "same 1,836 beats; 138 vs 1,794 bursts; 1.86x faster"],
            ["Runtime follows an exact law",
             "cycles = 3 x bursts + 2 x beats + compute"],
            ["Input-stationary is fastest on real hardware",
             "5,220 vs 5,808 cycles; the model agrees"],
            ["Prediction accuracy vs measured RTL",
             "traffic exact 17/19; cycles 0.84 % mean; 0 ranking inversions"]]
    table(s, rows, Inches(0.6), Inches(4.05), Inches(12.13), Inches(2.5),
          col_widths=[Inches(5.6), Inches(6.53)], size=14)

    caption(s, "All evidence: results/golden_check/ (measured RTL) and "
               "results/chooser/ (model). Reproducible in "
               "results/thesis_notebook/chooser_validation.ipynb.",
            Inches(6.7), size=13)


def main() -> None:
    prs = new_deck()
    slide_flow(prs)
    slide_provenance(prs)
    slide_why_no_big_rtl(prs)
    slide_verification_fixed(prs)
    slide_generalisation(prs)
    slide_summary(prs)
    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
